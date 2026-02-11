"""Translate PowerPoint files to Spanish using GPT-4o.

Walks all slides -> shapes -> text frames -> paragraphs -> runs.
Sends all text from a PPTX in one GPT-4o call with structural markers.
Writes translated text back, preserving all formatting.
Saves as {original_stem}_es.pptx.
"""

from __future__ import annotations

import json
import os
import re
import sys
from pathlib import Path

from lxml import etree

from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Emu

# XML namespace for DrawingML text elements (used in SmartArt diagrams)
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_DIAGRAM_DATA_RELTYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData"
)

load_dotenv()

BASE = Path(r"c:\Users\rcox\INSULATIONS, INC\Supervisory Training - Documents")
TRANSLATION_MODEL = "o3"

SYSTEM_PROMPT = """\
You are a professional translator specializing in English-to-Spanish translation \
for construction industry supervisor training presentations. The audience is field \
supervisors (foremen, superintendents, general foremen) at I&I Soft Craft Solutions, \
a company providing insulation, coatings, scaffolding, refractory, fireproofing, \
and heat tracing services.

Context & Tone:
- Use professional Latin American Spanish with the formal "usted" register, \
appropriate for corporate training materials.
- Tone should be warm, respectful, and encouraging — like an experienced mentor \
speaking to a valued team member. Avoid both cold corporate language and overly \
casual speech.
- Use first-person plural ("nosotros", "nuestro") to create a sense of shared \
ownership and belonging.

Slide Structure:
- Each text item has a "role" field: "title", "subtitle", or "body".
- TITLE text must be a short heading (2-5 words). Translate the meaning concisely — \
do NOT expand a title into a sentence or merge it with body content. \
Example: "Company Overview" → "Descripción de la Empresa", "Core Values" → "Valores Fundamentales".
- SUBTITLE text should also be brief (one short phrase).
- BODY text can be longer but must still fit within a PowerPoint text box. \
Keep translations concise — aim for equal or shorter length than the English original.

Translation Quality:
- Preserve full meaning. Condense phrasing to fit slide text boxes (Spanish is \
~15-20%% longer than English), but NEVER drop substantive content.
- Actively shorten wordy phrases. Prefer concise equivalents over literal translations \
when the literal version would be significantly longer.
- When the English is abstract or corporate, ground it with brief concrete phrasing \
relevant to a field supervisor's daily work.
- This is persuasive training content. Ensure the Spanish conveys why each point \
matters to the listener personally, not just what to do.
- Replace English idioms and colloquialisms with natural Spanish equivalents. \
Never translate idioms literally.

Language & Style:
- Use appropriate construction trade terminology in Spanish (capataz, superintendente, \
andamiaje, trazado térmico, ignifugación, etc.).
- When translating bulleted lists or enumerations, use parallel grammatical structure \
(all nouns, all infinitives, or all imperative — be consistent).
- Use short, punchy declarative sentences for emphasis on key messages and core values. \
Vary sentence length for rhetorical impact.
- Always use correct Spanish punctuation: inverted question/exclamation marks (¿ ¡), \
all accent marks (á, é, í, ó, ú, ñ), and proper em-dashes.

Terminology Glossary (use these exact terms):
- insulation → aislamiento
- coatings → recubrimientos (NOT revestimientos)
- scaffolding → andamiaje
- refractory → refractario
- fireproofing → ignifugación
- heat tracing → trazado térmico
- foreman → capataz
- superintendent → superintendente
- general foreman → capataz general
- safety → seguridad
- quality → calidad

Preservation Rules:
- Never translate the company name "I&I Soft Craft Solutions" or abbreviation "INI". \
Keep all branded service names in their original form.
- Preserve any numbers, acronyms, and proper nouns as-is.
- Preserve any formatting markers like bullet characters (•, -, etc.).
- Do NOT translate placeholder text or empty strings — return them as-is.
- Preserve line breaks (\\n) exactly as they appear in the original text.
- If a text string is ONLY whitespace, numbers, or punctuation, return it unchanged.
- Return valid JSON matching the schema exactly.
"""


def _shape_role(shape) -> str:
    """Determine the role of a shape: 'title', 'subtitle', or 'body'."""
    if shape.is_placeholder:
        ph_type = shape.placeholder_format.type
        if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            return "title"
        if ph_type == PP_PLACEHOLDER.SUBTITLE:
            return "subtitle"
    # Also detect by shape name convention
    name = shape.name.lower()
    if "title" in name:
        return "title"
    if "subtitle" in name:
        return "subtitle"
    return "body"


def collect_texts(prs: Presentation) -> list[dict]:
    """Collect all text runs from a presentation with location metadata.

    Returns a list of dicts with keys: id, text, role, location (for debugging).
    Skips empty/whitespace-only runs.
    """
    texts = []
    idx = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        # Slide body shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                role = _shape_role(shape)
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    for run_idx, run in enumerate(para.runs):
                        if run.text and run.text.strip():
                            texts.append({
                                "id": idx,
                                "text": run.text,
                                "role": role,
                                "location": f"slide {slide_num}, shape '{shape.name}', para {para_idx}, run {run_idx}",
                                "type": "slide",
                            })
                            idx += 1

            # Tables
            if shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        for para_idx, para in enumerate(cell.text_frame.paragraphs):
                            for run_idx, run in enumerate(para.runs):
                                if run.text and run.text.strip():
                                    texts.append({
                                        "id": idx,
                                        "text": run.text,
                                        "location": f"slide {slide_num}, table '{shape.name}', row {row_idx}, col {col_idx}, para {para_idx}, run {run_idx}",
                                        "type": "table",
                                    })
                                    idx += 1

        # SmartArt diagrams (stored as diagram data XML parts)
        slide_part = slide.part
        for rel in slide_part.rels.values():
            if rel.reltype == _DIAGRAM_DATA_RELTYPE:
                dgm_root = etree.fromstring(rel.target_part.blob)
                for t_idx, t_elem in enumerate(
                    dgm_root.iter(f"{{{_A_NS}}}t")
                ):
                    if t_elem.text and t_elem.text.strip():
                        texts.append({
                            "id": idx,
                            "text": t_elem.text,
                            "role": "body",
                            "location": f"slide {slide_num}, smartart, text {t_idx}",
                            "type": "smartart",
                            "_rel_rId": rel.rId,
                            "_t_index": t_idx,
                        })
                        idx += 1

        # Speaker notes
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            for para_idx, para in enumerate(notes_tf.paragraphs):
                for run_idx, run in enumerate(para.runs):
                    if run.text and run.text.strip():
                        texts.append({
                            "id": idx,
                            "text": run.text,
                            "location": f"slide {slide_num}, notes, para {para_idx}, run {run_idx}",
                            "type": "notes",
                        })
                        idx += 1

    return texts


def translate_texts(texts: list[dict], client: OpenAI) -> dict[int, str]:
    """Send all collected texts to GPT-4o for translation.

    Returns a mapping of id -> translated_text.
    Processes in batches of 80 to stay within token limits.
    """
    BATCH_SIZE = 80
    translations = {}

    for batch_start in range(0, len(texts), BATCH_SIZE):
        batch = texts[batch_start:batch_start + BATCH_SIZE]
        batch_translations = _translate_batch(batch, client)
        translations.update(batch_translations)

    return translations


def _translate_batch(texts: list[dict], client: OpenAI) -> dict[int, str]:
    """Translate a batch of texts via GPT-4o."""
    items = [{"id": t["id"], "text": t["text"], "role": t.get("role", "body")} for t in texts]
    user_message = json.dumps(items, ensure_ascii=False)

    schema = {
        "type": "object",
        "properties": {
            "translations": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "id": {"type": "integer"},
                        "translated_text": {"type": "string"},
                    },
                    "required": ["id", "translated_text"],
                    "additionalProperties": False,
                },
            },
        },
        "required": ["translations"],
        "additionalProperties": False,
    }

    response = client.responses.create(
        model=TRANSLATION_MODEL,
        input=[
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": user_message},
        ],
        text={
            "format": {
                "type": "json_schema",
                "name": "pptx_translation",
                "strict": True,
                "schema": schema,
            }
        },
    )

    parsed = json.loads(response.output_text)
    return {t["id"]: t["translated_text"] for t in parsed["translations"]}


def _restore_whitespace(original: str, translated: str) -> str:
    """Re-apply leading/trailing whitespace from original text to translation."""
    leading = original[: len(original) - len(original.lstrip())]
    trailing = original[len(original.rstrip()) :]
    return leading + translated.strip() + trailing


def apply_translations(prs: Presentation, texts: list[dict], translations: dict[int, str]) -> None:
    """Write translated text back into the presentation, preserving formatting."""
    # Build a lookup from text entries
    text_by_id = {t["id"]: t for t in texts}

    idx = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        # Slide body shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text and run.text.strip():
                            if idx in translations:
                                run.text = _restore_whitespace(
                                    run.text, translations[idx]
                                )
                            idx += 1

            # Tables
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.text and run.text.strip():
                                    if idx in translations:
                                        run.text = _restore_whitespace(
                                            run.text, translations[idx]
                                        )
                                    idx += 1

        # SmartArt diagrams
        slide_part = slide.part
        for rel in slide_part.rels.values():
            if rel.reltype == _DIAGRAM_DATA_RELTYPE:
                dgm_part = rel.target_part
                dgm_root = etree.fromstring(dgm_part.blob)
                modified = False
                for t_elem in dgm_root.iter(f"{{{_A_NS}}}t"):
                    if t_elem.text and t_elem.text.strip():
                        if idx in translations:
                            t_elem.text = _restore_whitespace(
                                t_elem.text, translations[idx]
                            )
                            modified = True
                        idx += 1
                if modified:
                    dgm_part._blob = etree.tostring(dgm_root, xml_declaration=True, encoding="UTF-8", standalone=True)

        # Speaker notes
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            for para in notes_tf.paragraphs:
                for run in para.runs:
                    if run.text and run.text.strip():
                        if idx in translations:
                            run.text = _restore_whitespace(
                                run.text, translations[idx]
                            )
                        idx += 1


def translate_pptx(pptx_path: Path, client: OpenAI) -> Path:
    """Translate a single PPTX file and save as _es.pptx."""
    output_path = pptx_path.parent / f"{pptx_path.stem}_es.pptx"

    print(f"  Loading: {pptx_path.name}")
    prs = Presentation(str(pptx_path))

    print(f"  Collecting text...")
    texts = collect_texts(prs)
    print(f"  Found {len(texts)} text runs to translate")

    if not texts:
        print(f"  No text found, saving copy as-is")
        prs.save(str(output_path))
        return output_path

    print(f"  Translating via GPT-4o ({len(texts)} strings)...")
    translations = translate_texts(texts, client)
    print(f"  Got {len(translations)} translations back")

    if len(translations) != len(texts):
        print(f"  WARNING: Expected {len(texts)} translations, got {len(translations)}!")
        missing = [t["id"] for t in texts if t["id"] not in translations]
        if missing:
            print(f"  Missing IDs: {missing[:10]}{'...' if len(missing) > 10 else ''}")

    print(f"  Applying translations...")
    apply_translations(prs, texts, translations)

    print(f"  Saving: {output_path.name}")
    prs.save(str(output_path))

    return output_path


def find_pptx_files() -> list[Path]:
    """Find all PPTX files to translate (skip already-translated _es files)."""
    files = []
    for pptx in BASE.rglob("*.pptx"):
        # Skip already-translated files
        if pptx.stem.endswith("_es"):
            continue
        # Skip temp/lock files
        if pptx.name.startswith("~"):
            continue
        files.append(pptx)

    files.sort(key=lambda p: p.name)
    return files


def main():
    api_key = os.environ.get("OPENAI_API_KEY", "").strip()
    if not api_key:
        print("ERROR: OPENAI_API_KEY not set. Add it to .env file.")
        sys.exit(1)

    client = OpenAI(api_key=api_key)

    pptx_files = find_pptx_files()
    print(f"Found {len(pptx_files)} PPTX files:\n")

    for f in pptx_files:
        es_path = f.parent / f"{f.stem}_es.pptx"
        exists = es_path.exists()
        print(f"  {'[DONE]' if exists else '[ .. ]'} {f.name}")

    pending = [f for f in pptx_files if not (f.parent / f"{f.stem}_es.pptx").exists()]
    print(f"\n{len(pending)} remaining to translate.\n")

    if not pending:
        print("All files already translated!")
        return

    for i, pptx_path in enumerate(pending, 1):
        print(f"\n{'='*60}")
        print(f"[{i}/{len(pending)}] {pptx_path.name}")
        print(f"{'='*60}")

        try:
            output = translate_pptx(pptx_path, client)
            size_kb = output.stat().st_size / 1024
            print(f"  OK: {output.name} ({size_kb:.0f} KB)")
        except Exception as exc:
            print(f"  ERROR: {exc}")
            import traceback
            traceback.print_exc()

    print(f"\n{'='*60}")
    print("TRANSLATION COMPLETE")
    print(f"{'='*60}")

    done = sum(1 for f in pptx_files if (f.parent / f"{f.stem}_es.pptx").exists())
    print(f"\n{done}/{len(pptx_files)} PPTX files translated.")


if __name__ == "__main__":
    main()
